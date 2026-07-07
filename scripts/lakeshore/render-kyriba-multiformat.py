#!/usr/bin/env python3
"""
Lakeshore Kyriba Move — multi-format artifact generator.

Takes the 8 board-grade HTML decks produced by the expert-kernel renderers and
emits the CORRECT downstream format per artifact type:

  charter / discover / mobilize  -> DOCX + PDF   (narrative briefs)
  solution-architecture / cfo    -> PPTX + PDF   (board decks)
  costed-business-case           -> PPTX + PDF + XLSX
  estimate-model                 -> XLSX + PDF   (financial model)
  master-dossier                 -> PDF + PPTX   (the assembled book)

HTML remains the interactive board-grade deck. These are the portable/editable
derivatives a CFO / SI / board expects. Content is extracted from the kernel
HTML (real kernel numbers + narrative), so no new numbers are introduced.
"""
from __future__ import annotations
import re, html as H
from pathlib import Path
from bs4 import BeautifulSoup

ROOT = Path(__file__).resolve().parents[2]
SRC = ROOT / "docs/build/lakeshore-enterprise-context/move-artifacts/board-grade"
OUT = SRC / "exports"
OUT.mkdir(parents=True, exist_ok=True)
WATERMARK = "SYNTHETIC — LAKESHORE PILOT — NOT REAL DATA — ABARVA AI"

# artifact id -> (label, [formats])
PLAN = {
    "lakeshore-kyriba-01-charter-skeleton":   ("Charter Skeleton",            ["docx", "pdf"]),
    "lakeshore-kyriba-02-discover-brief":     ("Discover Brief",              ["docx", "pdf"]),
    "lakeshore-kyriba-03-solution-architecture": ("Solution Architecture Pack", ["pptx", "pdf"]),
    "lakeshore-kyriba-04-costed-business-case": ("Costed Business Case",       ["pptx", "pdf", "xlsx"]),
    "lakeshore-kyriba-05-estimate-model":     ("Estimate & Financial Model",  ["xlsx", "pdf"]),
    "lakeshore-kyriba-06-cfo-pack":           ("CFO Pack",                    ["pptx", "pdf"]),
    "lakeshore-kyriba-07-mobilize-packet":    ("Mobilize & Go-Decision Packet", ["docx", "pdf"]),
    "lakeshore-kyriba-08-master-dossier":     ("Master Move Dossier",         ["pdf", "pptx"]),
}


def extract_sections(html: str):
    """Return [{title, paras:[str], tables:[[ [cell,...] ]]}] from a kernel deck."""
    soup = BeautifulSoup(html, "html.parser")
    deck_title = (soup.find(["h1"]) or soup.find(["title"]))
    deck_title = re.sub(r"\s+", " ", deck_title.get_text(" ", strip=True)) if deck_title else "Lakeshore Kyriba Move"
    sections = []
    nodes = soup.find_all("section")
    if not nodes:
        nodes = [soup]
    for sec in nodes:
        head = sec.find(["h1", "h2", "h3"])
        title = re.sub(r"\s+", " ", head.get_text(" ", strip=True)) if head else ""
        paras = []
        for p in sec.find_all(["p", "li"]):
            t = re.sub(r"\s+", " ", p.get_text(" ", strip=True))
            if t and len(t) > 1:
                paras.append(t)
        tables = []
        for tbl in sec.find_all("table"):
            rows = []
            for tr in tbl.find_all("tr"):
                cells = [re.sub(r"\s+", " ", c.get_text(" ", strip=True)) for c in tr.find_all(["th", "td"])]
                if any(cells):
                    rows.append(cells)
            if rows:
                tables.append(rows)
        # de-dup paragraphs, cap
        seen = set(); clean = []
        for t in paras:
            if t not in seen:
                seen.add(t); clean.append(t)
        sections.append({"title": title, "paras": clean[:18], "tables": tables})
    return deck_title, sections


# ---------------- DOCX ----------------
def emit_docx(path, title, sections):
    from docx import Document
    from docx.shared import Pt, RGBColor
    d = Document()
    cp = d.core_properties
    cp.title = title; cp.author = "AbarVa Expert Kernel"; cp.comments = WATERMARK
    hdr = d.sections[0].header.paragraphs[0]; hdr.text = WATERMARK
    hdr.runs[0].font.size = Pt(7); hdr.runs[0].font.color.rgb = RGBColor(0x8a, 0x93, 0xa6)
    d.add_heading(title, level=0)
    for s in sections:
        if s["title"]:
            d.add_heading(s["title"], level=1)
        for p in s["paras"]:
            d.add_paragraph(p)
        for tbl in s["tables"]:
            if not tbl:
                continue
            cols = max(len(r) for r in tbl)
            t = d.add_table(rows=0, cols=cols); t.style = "Light Grid Accent 1"
            for r in tbl:
                cells = t.add_row().cells
                for i, v in enumerate(r):
                    cells[i].text = v
    d.save(path)


# ---------------- PPTX ----------------
def emit_pptx(path, title, sections):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    INK = RGBColor(0x0c, 0x1a, 0x3a); MUTE = RGBColor(0x6b, 0x72, 0x80); NAVY = RGBColor(0x1F, 0x2A, 0x44)
    # cover
    c = prs.slides.add_slide(blank)
    tb = c.shapes.add_textbox(Inches(0.6), Inches(2.4), Inches(12), Inches(2.5)).text_frame
    tb.word_wrap = True
    r = tb.paragraphs[0].add_run(); r.text = title; r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = INK
    sub = tb.add_paragraph().add_run(); sub.text = WATERMARK; sub.font.size = Pt(11); sub.font.color.rgb = MUTE
    for s in sections:
        sl = prs.slides.add_slide(blank)
        head = sl.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.9)).text_frame
        head.word_wrap = True
        hr = head.paragraphs[0].add_run(); hr.text = s["title"] or title
        hr.font.size = Pt(22); hr.font.bold = True; hr.font.color.rgb = NAVY
        body = sl.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.6)).text_frame
        body.word_wrap = True
        bullets = s["paras"][:8]
        for i, p in enumerate(bullets):
            para = body.paragraphs[0] if i == 0 else body.add_paragraph()
            run = para.add_run(); run.text = "• " + (p[:240]); run.font.size = Pt(13); run.font.color.rgb = INK
        if s["tables"]:
            tbl = s["tables"][0][:8]
            rows = len(tbl); cols = max(len(r) for r in tbl)
            if rows and cols:
                gt = sl.shapes.add_table(rows, cols, Inches(0.6), Inches(4.4), Inches(12.1), Inches(2.6)).table
                for ri, row in enumerate(tbl):
                    for ci in range(cols):
                        gt.cell(ri, ci).text = row[ci] if ci < len(row) else ""
    prs.save(path)


# ---------------- XLSX ----------------
def emit_xlsx(path, title, sections):
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill
    wb = Workbook(); ws0 = wb.active; ws0.title = "Overview"
    ws0.append([WATERMARK]); ws0.append([title]); ws0.append([])
    ws0["A1"].font = Font(bold=True, color="8a93a6", size=9)
    ws0["A2"].font = Font(bold=True, size=13)
    tnum = 0
    for s in sections:
        if s["paras"]:
            for p in s["paras"][:4]:
                ws0.append([s["title"], p])
        for tbl in s["tables"]:
            tnum += 1
            ws = wb.create_sheet(f"T{tnum}_{(s['title'] or 'table')[:22]}".replace("/", "-")[:31] or f"T{tnum}")
            ws.append([WATERMARK])
            ws["A1"].font = Font(bold=True, color="8a93a6", size=9)
            for ri, row in enumerate(tbl):
                ws.append(row)
                if ri == 0:
                    for ci in range(len(row)):
                        ws.cell(row=2, column=ci + 1).font = Font(bold=True)
                        ws.cell(row=2, column=ci + 1).fill = PatternFill("solid", fgColor="D8DEE9")
    if tnum == 0:
        ws0.append([]); ws0.append(["(No tabular exhibits in source deck.)"])
    wb.save(path)


# ---------------- PDF ----------------
def emit_pdf(path, title, sections):
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, ListFlowable, ListItem)
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    styles = getSampleStyleSheet()
    h1 = ParagraphStyle("h1", parent=styles["Heading1"], textColor=colors.HexColor("#0c1a3a"), fontSize=16)
    h2 = ParagraphStyle("h2", parent=styles["Heading2"], textColor=colors.HexColor("#1F2A44"), fontSize=12)
    body = ParagraphStyle("body", parent=styles["BodyText"], fontSize=9.3, leading=12.5)
    eyebrow = ParagraphStyle("eb", parent=styles["BodyText"], fontSize=7.5, textColor=colors.HexColor("#6b7280"))

    def wm(canvas, doc):
        canvas.saveState(); canvas.setFont("Helvetica", 7); canvas.setFillColor(colors.HexColor("#8a93a6"))
        canvas.drawString(0.6 * inch, 0.4 * inch, WATERMARK)
        canvas.drawRightString(7.9 * inch, 0.4 * inch, f"page {doc.page}"); canvas.restoreState()

    doc = SimpleDocTemplate(str(path), pagesize=letter, topMargin=0.7 * inch, bottomMargin=0.7 * inch,
                            leftMargin=0.7 * inch, rightMargin=0.7 * inch, title=title)
    flow = [Paragraph(WATERMARK, eyebrow), Paragraph(title, h1), Spacer(1, 0.15 * inch)]
    for s in sections:
        if s["title"]:
            flow.append(Paragraph(H.escape(s["title"]), h2))
        if s["paras"]:
            items = [ListItem(Paragraph(H.escape(p), body), leftIndent=8) for p in s["paras"][:12]]
            flow.append(ListFlowable(items, bulletType="bullet")); flow.append(Spacer(1, 0.05 * inch))
        for tbl in s["tables"]:
            data = [[H.escape(c)[:60] for c in (row + [""] * (max(len(r) for r in tbl) - len(row)))] for row in tbl[:14]]
            t = Table(data, repeatRows=1, hAlign="LEFT")
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F2A44")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTSIZE", (0, 0), (-1, -1), 7),
                ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#c9ccd6")),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f3f4f7")]),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]))
            flow.append(t); flow.append(Spacer(1, 0.1 * inch))
        flow.append(Spacer(1, 0.08 * inch))
    doc.build(flow, onFirstPage=wm, onLaterPages=wm)


EMIT = {"docx": emit_docx, "pptx": emit_pptx, "xlsx": emit_xlsx, "pdf": emit_pdf}


def main():
    rows = []
    for aid, (label, formats) in PLAN.items():
        html_path = SRC / f"{aid}.html"
        if not html_path.exists():
            print("MISSING", html_path); continue
        title, sections = extract_sections(html_path.read_text(encoding="utf-8"))
        for fmt in formats:
            out_path = OUT / f"{aid}.{fmt}"
            EMIT[fmt](out_path, title, sections)
            kb = out_path.stat().st_size / 1024
            rows.append((aid, label, fmt, f"{kb:.0f} KB"))
            print(f"{label:32s} {fmt.upper():4s} -> {out_path.name} [{kb:.0f}kb]")
    # index
    idx = ["# Lakeshore Kyriba — Multi-format Move Artifacts", "",
           "Correct downstream format per artifact type (HTML decks remain the interactive board-grade source).", "",
           "| Artifact | Format | File | Size |", "|---|---|---|---|"]
    for aid, label, fmt, sz in rows:
        idx.append(f"| {label} | {fmt.upper()} | `exports/{aid}.{fmt}` | {sz} |")
    (OUT / "INDEX.md").write_text("\n".join(idx) + "\n", encoding="utf-8")
    print(f"\nWrote {len(rows)} files to {OUT}")


if __name__ == "__main__":
    main()
